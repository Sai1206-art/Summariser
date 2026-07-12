#!/usr/bin/env python3
"""
Zapplon — Services & Portfolio pitch deck generator (16:9).

Brand system:
  bg charcoal-black #14110D · amber #F3A81E · cream #F4ECDF · tan #A99B85
  One accent color. Big bold headlines, light body. Trebuchet MS.

Rules honored:
  - No invented metrics: every unknown figure stays as a [[ token ]].
  - Charts are native editable pptx chart objects; their values are
    illustrative shapes only and are labeled as such on-slide.
  - One consistent template across all slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import copy

# ---------------------------------------------------------------- brand
BG     = RGBColor(0x14, 0x11, 0x0D)
AMBER  = RGBColor(0xF3, 0xA8, 0x1E)
CREAM  = RGBColor(0xF4, 0xEC, 0xDF)
TAN    = RGBColor(0xA9, 0x9B, 0x85)
CARD   = RGBColor(0x1E, 0x1A, 0x14)   # slightly lifted panel on charcoal
CARD2  = RGBColor(0x26, 0x21, 0x19)   # chip fill
FONT   = "Trebuchet MS"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

TOKENS = []          # every [[ token ]] used, in order, for the fill-in list
def tok(name):
    if name not in TOKENS:
        TOKENS.append(name)
    return f"[[ {name} ]]"

# ---------------------------------------------------------------- helpers
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def txt(s, l, t, w, h, runs, size=14, color=CREAM, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.0, wrap=True):
    """runs: str, or list of paragraphs; each paragraph: str or list of (text, dict) runs."""
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(para, str):
            para = [(para, {})]
        for text, style in para:
            r = p.add_run(); r.text = text
            f = r.font
            f.name = style.get("font", FONT)
            f.size = Pt(style.get("size", size))
            f.bold = style.get("bold", bold)
            f.color.rgb = style.get("color", color)
            if style.get("spacing_pt"):
                pass
    return box

def panel(s, l, t, w, h, fill=CARD, line=None, radius=0.08):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def kicker(s, text):
    txt(s, Inches(0.55), Inches(0.38), Inches(9), Inches(0.35),
        [[(text.upper(), {"size": 11, "bold": True, "color": AMBER})]])
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.57), Inches(0.72),
                            Inches(0.45), Pt(2.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = AMBER
    ln.line.fill.background(); ln.shadow.inherit = False

def headline(s, text, size=32, top=Inches(0.86), left=Inches(0.55), width=Inches(11.5)):
    txt(s, left, top, width, Inches(1.0),
        [[(text, {"size": size, "bold": True, "color": CREAM})]])

def footer(s, page):
    txt(s, Inches(0.55), Inches(7.08), Inches(6), Inches(0.3),
        [[("Zapplon", {"size": 9, "bold": True, "color": TAN}),
          ("  ·  outcomes, not hours", {"size": 9, "color": TAN})]])
    txt(s, Inches(12.35), Inches(7.08), Inches(0.7), Inches(0.3),
        [[(f"{page:02d}", {"size": 9, "bold": True, "color": TAN})]],
        align=PP_ALIGN.RIGHT)

def logo(s, l=Inches(11.35), t=Inches(0.35), small=True):
    """Amber flame mark + wordmark."""
    flame = s.shapes.add_shape(MSO_SHAPE.TEAR, l, t, Inches(0.26), Inches(0.26))
    flame.rotation = 315
    flame.fill.solid(); flame.fill.fore_color.rgb = AMBER
    flame.line.fill.background(); flame.shadow.inherit = False
    txt(s, l + Inches(0.32), t - Inches(0.05), Inches(1.15), Inches(0.35),
        [[("Zapplon", {"size": 14, "bold": True, "color": CREAM})]], wrap=False)

def chip(s, l, t, w, big, small_lbl, h=Inches(0.95)):
    panel(s, l, t, w, h, fill=CARD2)
    txt(s, l + Inches(0.14), t + Inches(0.10), w - Inches(0.28), Inches(0.45),
        [[(big, {"size": 13, "bold": True, "color": AMBER})]])
    txt(s, l + Inches(0.14), t + Inches(0.52), w - Inches(0.28), Inches(0.38),
        [[(small_lbl.upper(), {"size": 8, "bold": True, "color": TAN})]])

def browser_frame(s, l, t, w, h, url, note_lines):
    """Clean browser frame with a labeled drop-zone for the real screenshot."""
    panel(s, l, t, w, h, fill=CARD, line=RGBColor(0x3A, 0x33, 0x28), radius=0.045)
    bar_h = Inches(0.42)
    bar = panel(s, l, t, w, bar_h, fill=CARD2, radius=0.045)
    for i in range(3):
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.18 + i * 0.22),
                               t + Inches(0.15), Inches(0.11), Inches(0.11))
        d.fill.solid(); d.fill.fore_color.rgb = TAN
        d.line.fill.background(); d.shadow.inherit = False
    urlbox = panel(s, l + Inches(0.9), t + Inches(0.08), w - Inches(1.2),
                   Inches(0.26), fill=BG, radius=0.5)
    txt(s, l + Inches(1.05), t + Inches(0.085), w - Inches(1.4), Inches(0.26),
        [[(url, {"size": 9, "color": TAN})]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, l + Inches(0.4), t + bar_h + Inches(0.35), w - Inches(0.8),
        h - bar_h - Inches(0.7),
        [[("DROP SCREENSHOT HERE", {"size": 13, "bold": True, "color": AMBER})]]
        + [[(n, {"size": 10, "color": TAN})] for n in note_lines],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.25)

def style_chart(ch, series_amber=True):
    ch.has_legend = False
    ch.font.name = FONT
    ch.font.size = Pt(8)
    ch.font.color.rgb = TAN
    for plot in ch.plots:
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = AMBER if (series_amber and i == 0) else TAN
            ser.format.line.fill.background()
    for ax_name in ("category_axis", "value_axis"):
        try:
            ax = getattr(ch, ax_name)
            ax.format.line.color.rgb = RGBColor(0x3A, 0x33, 0x28)
            ax.tick_labels.font.size = Pt(8)
            ax.tick_labels.font.color.rgb = TAN
            ax.has_major_gridlines = False
            ax.has_minor_gridlines = False
        except Exception:
            pass

def mini_bar_chart(s, l, t, w, h, cats, vals, title_note):
    chart_w = Inches(2.6)
    data = CategoryChartData()
    data.categories = cats
    data.add_series("value", vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, chart_w, h, data)
    ch = gf.chart
    ch.has_title = False
    style_chart(ch)
    txt(s, l + chart_w + Inches(0.15), t, w - chart_w - Inches(0.15), h,
        [[(title_note, {"size": 8, "color": TAN})]],
        anchor=MSO_ANCHOR.MIDDLE, spacing=1.15)
    return ch

# =================================================================
# SLIDE 0 — COVER
# =================================================================
s = slide()
# amber flame mark, large
flame = s.shapes.add_shape(MSO_SHAPE.TEAR, Inches(0.62), Inches(2.05), Inches(0.55), Inches(0.55))
flame.rotation = 315
flame.fill.solid(); flame.fill.fore_color.rgb = AMBER
flame.line.fill.background(); flame.shadow.inherit = False
txt(s, Inches(0.55), Inches(2.75), Inches(9), Inches(1.0),
    [[("Zapplon", {"size": 54, "bold": True, "color": CREAM})]])
txt(s, Inches(0.6), Inches(3.85), Inches(10.5), Inches(0.6),
    [[("Outcomes, not hours.", {"size": 22, "bold": True, "color": AMBER})]])
txt(s, Inches(0.6), Inches(4.5), Inches(9.5), Inches(0.8),
    [[("Performance growth marketing for US local & consumer businesses — "
       "services & portfolio.", {"size": 13, "color": TAN})]], spacing=1.2)
txt(s, Inches(0.6), Inches(6.35), Inches(8), Inches(0.35),
    [[("zapplon.com", {"size": 11, "bold": True, "color": CREAM})]])
footer(s, 1)

# =================================================================
# SLIDE 1 — OUR SERVICES
# =================================================================
s = slide()
kicker(s, "Services")
headline(s, "What we do.")
logo(s)

services = [
    ("◎", "Local SEO & Google Business Profile", "Own the map pack where buying decisions start."),
    ("▣", "Website design & build", "Fast, conversion-first sites that turn visits into bookings."),
    ("▶", "Paid ads", "Google, Meta & delivery-app campaigns tuned to cost-per-lead."),
    ("✎", "Content & social", "A steady drumbeat that keeps your brand chosen."),
    ("✉", "Email & cold outreach", "Pipelines that warm, nurture and book."),
    ("⚙", "AI automation", "Receptionists, follow-ups and workflows that never sleep."),
]
cw, chh = Inches(4.0), Inches(1.55)
gx, gy = Inches(0.55), Inches(1.75)
for i, (icon, name, desc) in enumerate(services):
    col, row = i % 3, i // 3
    l = gx + col * (cw + Inches(0.12))
    t = gy + row * (chh + Inches(0.14))
    panel(s, l, t, cw, chh, fill=CARD)
    txt(s, l + Inches(0.2), t + Inches(0.16), Inches(0.5), Inches(0.5),
        [[(icon, {"size": 18, "color": AMBER})]])
    txt(s, l + Inches(0.72), t + Inches(0.18), cw - Inches(0.9), Inches(0.45),
        [[(name, {"size": 12.5, "bold": True, "color": CREAM})]])
    txt(s, l + Inches(0.72), t + Inches(0.62), cw - Inches(0.9), Inches(0.8),
        [[(desc, {"size": 10, "color": TAN})]], spacing=1.1)

# bottom metrics band
band_t = Inches(5.35)
panel(s, Inches(0.55), band_t, Inches(12.24), Inches(1.35), fill=CARD2)
band = [
    (tok("clients_served"), "clients"),
    (tok("campaigns_managed"), "campaigns"),
    (tok("leads_generated"), "leads"),
    (tok("websites_built"), "websites"),
    (tok("ai_automations"), "AI automations deployed"),
]
bw = Inches(12.24) / 5
for i, (big, lbl) in enumerate(band):
    l = Inches(0.55) + i * bw
    txt(s, l, band_t + Inches(0.22), bw, Inches(0.5),
        [[(big, {"size": 15, "bold": True, "color": AMBER})]], align=PP_ALIGN.CENTER)
    txt(s, l, band_t + Inches(0.78), bw, Inches(0.4),
        [[(lbl.upper(), {"size": 8.5, "bold": True, "color": TAN})]], align=PP_ALIGN.CENTER)
footer(s, 2)

# =================================================================
# PORTFOLIO TEMPLATE — one slide per business
# =================================================================
def portfolio_slide(page, vertical, city, screenshot_url, screenshot_note,
                    what_we_did, chart_note):
    s = slide()
    kicker(s, f"Portfolio · {vertical} · {city}")
    logo(s)

    # client logo placeholder + client name
    lg = panel(s, Inches(0.55), Inches(0.95), Inches(0.85), Inches(0.85),
               fill=CARD2, line=RGBColor(0x3A, 0x33, 0x28))
    txt(s, Inches(0.55), Inches(1.14), Inches(0.85), Inches(0.5),
        [[("LOGO", {"size": 8, "bold": True, "color": TAN})]], align=PP_ALIGN.CENTER)
    txt(s, Inches(1.55), Inches(0.98), Inches(5.4), Inches(0.55),
        [[(tok("client_name"), {"size": 24, "bold": True, "color": CREAM})]])
    txt(s, Inches(1.55), Inches(1.5), Inches(5.4), Inches(0.35),
        [[(f"{vertical} · {city}", {"size": 11, "color": TAN})]])

    # right ~55%: screenshot in browser frame
    browser_frame(s, Inches(6.35), Inches(0.95), Inches(6.45), Inches(5.75),
                  screenshot_url, screenshot_note)

    # left: what we did
    txt(s, Inches(0.55), Inches(2.1), Inches(5.4), Inches(0.35),
        [[("WHAT WE DID", {"size": 10, "bold": True, "color": AMBER})]])
    bullets = [[("—  ", {"size": 11, "color": AMBER}),
                (b, {"size": 11, "color": CREAM})] for b in what_we_did]
    txt(s, Inches(0.55), Inches(2.45), Inches(5.5), Inches(1.35), bullets, spacing=1.35)

    # results row — hero
    txt(s, Inches(0.55), Inches(3.85), Inches(5.4), Inches(0.35),
        [[("RESULTS", {"size": 10, "bold": True, "color": AMBER})]])
    cwid = Inches(1.78)
    chips = [
        (f'{tok("rank_before")} → {tok("rank_after")}', f'ranking · "{tok("keyword")}"'),
        (f'{tok("reviews_before")} → {tok("reviews_after")} ({tok("rating")}★)', "reviews"),
        (f'+{tok("traffic_growth_pct")}%', f'traffic · {tok("period")}'),
    ]
    for i, (big, lbl) in enumerate(chips):
        chip(s, Inches(0.55) + i * (cwid + Inches(0.1)), Inches(4.2), cwid, big, lbl)
    chips2 = [
        (f'{tok("leads_per_month")}/mo', "leads"),
        (f'{tok("conversion_pct")}%', "conversions"),
    ]
    for i, (big, lbl) in enumerate(chips2):
        chip(s, Inches(0.55) + i * (cwid + Inches(0.1)), Inches(5.25), cwid, big, lbl)
    chip(s, Inches(0.55) + 2 * (cwid + Inches(0.1)), Inches(5.25), cwid,
         tok("automation"), f'AI · {tok("volume")}/mo handled')

    # small native chart (editable) — illustrative shape only
    mini_bar_chart(s, Inches(0.55), Inches(6.28), Inches(5.54), Inches(0.66),
                   ["before", "after"], [1, 3], chart_note)
    footer(s, page)
    return s

portfolio_slide(
    3, "Dental", "San Jose",
    "google.com/search?q=best+dental+centers+in+san+jose",
    ["Google Places screenshot — “best dental centers in san jose”",
     "(attached asset: dental-san-jose search results)"],
    ["Google Business Profile optimization + local SEO",
     "Review engine — steady 5★ velocity",
     "Website refresh + booking funnel"],
    "editable chart — replace with real before/after values")

portfolio_slide(
    4, "Legal", "San Jose",
    "google.com/search?q=best+law+firms+in+san+jose",
    ["Google Businesses screenshot — “best law firms in san jose”",
     "(attached asset: legal-san-jose search results)"],
    ["Local SEO + practice-area landing pages",
     "Google Ads on high-intent defense keywords",
     "AI intake receptionist — 24/7 call answering"],
    "editable chart — replace with real before/after values")

portfolio_slide(
    5, "Fitness", "San Jose",
    "google.com/search?q=best+gyms+in+san+jose",
    ["Search / listing screenshot for the fitness client",
     "(insert client’s map-pack or listing capture)"],
    ["Google Business Profile + local SEO",
     "Meta ads — trial-membership offers",
     "Email win-back + referral automation"],
    "editable chart — replace with real before/after values")

portfolio_slide(
    6, "Martial arts", "San Jose",
    "google.com/search?q=best+martial+arts+in+san+jose",
    ["Search / listing screenshot for the martial-arts client",
     "(insert client’s map-pack or listing capture)"],
    ["Local SEO + review engine",
     "Paid ads — kids’ program enrollment",
     "AI follow-up on trial-class leads"],
    "editable chart — replace with real before/after values")

# =================================================================
# DESIGN BENCHMARKS — Posha & Corn Revolution (attributed, not ours)
# =================================================================
def benchmark_slide(page, name, owner_line, url, takeaways):
    s = slide()
    kicker(s, "Design benchmarks")
    headline(s, "The standard we build to.", size=28)
    logo(s)
    txt(s, Inches(0.55), Inches(1.55), Inches(5.6), Inches(0.9),
        [[("Not our build — a benchmark we admire. ", {"size": 11, "bold": True, "color": AMBER}),
          (owner_line, {"size": 11, "color": TAN})]], spacing=1.25)

    browser_frame(s, Inches(6.35), Inches(0.95), Inches(6.45), Inches(5.75),
                  url, [f"Hero screenshot — {name}", f"(clickable → {url})"])
    # make the frame area link out
    link = s.shapes.add_textbox(Inches(6.35), Inches(6.35), Inches(6.45), Inches(0.3))
    p = link.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"↗ {url}"
    r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = AMBER
    r.hyperlink.address = f"https://{url}"

    txt(s, Inches(0.55), Inches(2.6), Inches(5.4), Inches(0.35),
        [[("WHY IT SETS THE BAR", {"size": 10, "bold": True, "color": AMBER})]])
    txt(s, Inches(0.55), Inches(2.95), Inches(5.5), Inches(2.6),
        [[("—  ", {"size": 11, "color": AMBER}), (b, {"size": 11, "color": CREAM})]
         for b in takeaways], spacing=1.4)
    txt(s, Inches(0.55), Inches(5.9), Inches(5.5), Inches(0.8),
        [[("This is the calibre of craft we design toward on every Zapplon build.",
           {"size": 11, "bold": True, "color": CREAM})]], spacing=1.2)
    footer(s, page)

benchmark_slide(
    7, "Posha", "Posha — posha.com (all design credit to its owner).",
    "posha.com",
    ["One idea per screen — “See dinner cook itself.”",
     "Product photography does the selling; copy stays out of the way",
     "A single high-contrast CTA (“Buy now”) always in reach",
     "Warm, human art direction for a hardware product"])

benchmark_slide(
    8, "Corn Revolution", "Corn Revolution — built by RESN (cornrevolution.resn.global).",
    "cornrevolution.resn.global",
    ["Cinematic WebGL storytelling with restraint — dark canvas, one accent",
     "Data made emotional: germplasm history told as a living visual",
     "Typography as hero — massive, confident headlines",
     "Every interaction rewards curiosity without hurting performance"])

# =================================================================
# CLOSING — CTA
# =================================================================
s = slide()
kicker(s, "Next step")
flame = s.shapes.add_shape(MSO_SHAPE.TEAR, Inches(0.62), Inches(1.9), Inches(0.5), Inches(0.5))
flame.rotation = 315
flame.fill.solid(); flame.fill.fore_color.rgb = AMBER
flame.line.fill.background(); flame.shadow.inherit = False
txt(s, Inches(0.55), Inches(2.55), Inches(11.8), Inches(1.6),
    [[("Ready to be found,", {"size": 40, "bold": True, "color": CREAM})],
     [("chosen and booked?", {"size": 40, "bold": True, "color": AMBER})]],
    spacing=1.05)
contacts = [("zapplon.com", Inches(0.55)),
            (tok("email"), Inches(3.35)),
            (tok("phone"), Inches(6.15))]
for label, l in contacts:
    panel(s, l, Inches(4.75), Inches(2.6), Inches(0.6), fill=CARD2)
    txt(s, l, Inches(4.75), Inches(2.6), Inches(0.6),
        [[(label, {"size": 12, "bold": True, "color": CREAM})]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.55), Inches(5.75), Inches(11), Inches(0.4),
    [[("Zapplon · performance growth marketing · outcomes, not hours",
       {"size": 10, "color": TAN})]])
footer(s, 9)

# =================================================================
# FILL-IN LIST — every token used
# =================================================================
s = slide()
kicker(s, "Internal — remove before sending")
headline(s, "Fill-in list: replace every [[ token ]] with real data.", size=22)
logo(s)
notes = {
    "clients_served": "Slide 2 metrics band",
    "campaigns_managed": "Slide 2 metrics band",
    "leads_generated": "Slide 2 metrics band",
    "websites_built": "Slide 2 metrics band",
    "ai_automations": "Slide 2 metrics band",
    "client_name": "Each portfolio slide — real client name + logo",
    "rank_before": "Portfolio results chips", "rank_after": "Portfolio results chips",
    "keyword": "Ranked keyword per client",
    "reviews_before": "Portfolio results chips", "reviews_after": "Portfolio results chips",
    "rating": "Star rating per client",
    "traffic_growth_pct": "Traffic growth %", "period": "Measurement window",
    "leads_per_month": "Leads/mo per client", "conversion_pct": "Conversion % per client",
    "automation": "AI automation name (e.g. AI receptionist)",
    "volume": "Monthly volume the automation handles",
    "email": "Closing slide contact chip", "phone": "Closing slide contact chip",
}
col_w = Inches(6.1)
half = (len(TOKENS) + 1) // 2
for i, t_name in enumerate(TOKENS):
    col = 0 if i < half else 1
    row = i if i < half else i - half
    l = Inches(0.55) + col * (col_w + Inches(0.15))
    top = Inches(1.85) + row * Inches(0.46)
    txt(s, l, top, Inches(2.4), Inches(0.4),
        [[(f"[[ {t_name} ]]", {"size": 10.5, "bold": True, "color": AMBER})]])
    txt(s, l + Inches(2.45), top, col_w - Inches(2.45), Inches(0.4),
        [[(notes.get(t_name, ""), {"size": 10, "color": TAN})]])
txt(s, Inches(0.55), Inches(6.55), Inches(12), Inches(0.45),
    [[("Also insert: client logos, real listing screenshots into each browser frame, "
       "and real values into each editable chart (current chart values are illustrative only).",
       {"size": 10, "color": CREAM})]], spacing=1.2)
footer(s, 10)

OUT = "Zapplon-Services-Portfolio-Deck.pptx"
prs.save(OUT)
print(f"Saved {OUT} — {len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides, tokens: {len(TOKENS)}")
